#!/usr/bin/env python3
"""Build the lake trout PAV + methylation + annotation ecotype deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---------- palette (Ocean Gradient: deep water -> shallow) ----------
DEEP    = RGBColor(0x06, 0x2A, 0x42)   # near-black deep water (title bg)
MID     = RGBColor(0x06, 0x5A, 0x82)   # deep blue
TEAL    = RGBColor(0x1C, 0x72, 0x93)   # teal
SEAFOAM = RGBColor(0x2A, 0x9D, 0x8F)   # green accent (methylation)
CORAL   = RGBColor(0xE9, 0x96, 0x4E)   # warm accent (PAV / contrast)
GOLD    = RGBColor(0xC9, 0x7B, 0x2E)   # deeper warm (annotation / candidates)
ICE     = RGBColor(0xCA, 0xDC, 0xFC)   # ice blue light
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1A, 0x2B, 0x35)   # body text on light
MUTE    = RGBColor(0x5B, 0x70, 0x7B)   # muted caption
LIGHTBG = RGBColor(0xF4, 0xF7, 0xF9)   # light content bg

HEAD_FONT = "Georgia"
BODY_FONT = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        # lightweight soft shadow via XML
        from pptx.oxml.ns import qn
        efflst = el.makeelement(qn('a:effectLst'), {})
        outer = el.makeelement(qn('a:outerShdw'),
                               {'blurRad': '60000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val': '0A2230'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '28000'})
        clr.append(alpha)
        outer.append(clr)
        efflst.append(outer)
        el.append(efflst)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True, space_after=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        if isinstance(para, tuple):
            para = [para]
        for text, opt in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opt.get("font", BODY_FONT)
            f.size = Pt(opt.get("size", 16))
            f.bold = opt.get("bold", False)
            f.italic = opt.get("italic", False)
            f.color.rgb = opt.get("color", INK)
            if opt.get("spacing"):
                _set_spacing(r, opt["spacing"])
    return tb


def _set_spacing(run, pts):
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def bullets(s, x, y, w, h, items, size=16, color=INK, gap=10, bullet_color=TEAL):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        lead = "—  "
        r = p.add_run(); r.text = lead
        r.font.name = BODY_FONT; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = bullet_color
        if isinstance(item, str):
            item = [(item, {})]
        for text, opt in item:
            r = p.add_run(); r.text = text
            r.font.name = BODY_FONT
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", False)
            r.font.italic = opt.get("italic", False)
            r.font.color.rgb = opt.get("color", color)
    return tb


def kicker(s, x, y, text, color=CORAL):
    txt(s, x, y, 8, 0.35, [[(text.upper(), {"font": BODY_FONT, "size": 13, "bold": True,
                                            "color": color, "spacing": 2.5})]])


def page_no(s, n):
    txt(s, 12.4, 7.05, 0.7, 0.3, [[(f"{n:02d}", {"size": 11, "color": MUTE})]],
        align=PP_ALIGN.RIGHT)


N = 0
def num():
    global N
    N += 1
    return N

# ============================================================ 1. TITLE
s = slide(DEEP)
# layered depth bands
rect(s, 0, 5.0, 13.333, 2.5, MID)
rect(s, 0, 6.0, 13.333, 1.5, TEAL)
rect(s, 0, 6.8, 13.333, 0.7, SEAFOAM)
kicker(s, 0.9, 0.95, "Salvelinus namaycush  ·  Roberts Lab, UW SAFS", ICE)
txt(s, 0.9, 1.45, 11.5, 2.4, [
    [("From Variation to", {"font": HEAD_FONT, "size": 48, "bold": True, "color": WHITE})],
    [("Candidate Genes", {"font": HEAD_FONT, "size": 48, "bold": True, "color": ICE})],
], space_after=2)
txt(s, 0.92, 3.7, 11.5, 0.9,
    [[("Methylation, structural variation, and functional annotation of ",
       {"size": 20, "color": ICE}),
      ("lean vs. siscowet", {"size": 20, "bold": True, "color": WHITE}),
      (" lake trout ecotypes", {"size": 20, "color": ICE})]])
txt(s, 0.92, 6.55, 11.5, 0.4,
    [[("Rick Goetz · Sam White · Cristian Gallardo-Escarate · Steven Roberts",
       {"size": 13.5, "bold": True, "color": WHITE})]])
txt(s, 0.92, 7.0, 8, 0.4, [[("PacBio HiFi comparative epigenomics  ·  research talk",
                            {"size": 12, "color": ICE})]])

# ============================================================ 2. THE QUESTION
s = slide(WHITE)
kicker(s, 0.9, 0.6, "The biological question")
txt(s, 0.9, 1.0, 11.5, 1.0,
    [[("What makes a lean a lean and a siscowet a siscowet?",
       {"font": HEAD_FONT, "size": 32, "bold": True, "color": INK})]])

# two ecotype cards
def ecocard(x, color, name, depth, traits):
    rect(s, x, 2.15, 5.4, 3.55, LIGHTBG, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.15, 5.4, 0.85, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.55, 5.4, 0.45, color)  # square off bottom of header
    txt(s, x+0.4, 2.32, 4.6, 0.5, [[(name, {"font": HEAD_FONT, "size": 24, "bold": True, "color": WHITE})]])
    txt(s, x+0.4, 3.15, 4.6, 0.4, [[(depth, {"size": 15, "bold": True, "italic": True, "color": color})]])
    bullets(s, x+0.4, 3.7, 4.7, 1.9, traits, size=15, gap=9, bullet_color=color)

ecocard(0.9, TEAL, "Lean", "Shallow · nearshore",
        ["Streamlined, low lipid", "Forages shallow water", "Reference-like body plan"])
ecocard(7.05, MID, "Siscowet", "Deep · offshore",
        ["High lipid / fat stores", "Adapted to depth & pressure", "Robust, distinct morphology"])
txt(s, 0.9, 5.95, 11.5, 1.0,
    [[("Same species, same reference genome. ", {"size": 17, "color": INK}),
      ("Which genes carry the difference — and what phenotypes might they shape?",
       {"size": 17, "bold": True, "color": CORAL})]])
page_no(s, num())

# ============================================================ 3. THREE-LAYER FRAMEWORK
s = slide(DEEP)
kicker(s, 0.9, 0.55, "The framework", ICE)
txt(s, 0.9, 0.95, 11.5, 0.8, [[("Three layers, one shortlist of genes",
                                {"font": HEAD_FONT, "size": 31, "bold": True, "color": WHITE})]])

def layercard(x, tag, tag_c, title, body_runs):
    rect(s, x, 2.0, 3.7, 3.9, MID, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x+0.32, 2.3, 3.1, 0.4, [[(tag, {"size": 13.5, "bold": True, "color": tag_c, "spacing": 1.3})]])
    txt(s, x+0.32, 2.72, 3.1, 0.6, [[(title, {"font": HEAD_FONT, "size": 21, "bold": True, "color": WHITE})]])
    txt(s, x+0.32, 3.55, 3.1, 2.2, body_runs)

layercard(0.9, "LAYER 1 · PAV", CORAL, "The hardware",
          [[("Presence-absence variation — DNA segments ", {"size": 14.5, "color": ICE}),
            ("gained or lost", {"size": 14.5, "bold": True, "color": WHITE}),
            (".", {"size": 14.5, "color": ICE})],
           [("", {"size": 6})],
           [("Changes ", {"size": 14.5, "color": ICE}),
            ("what genes exist.", {"size": 14.5, "bold": True, "color": WHITE})]])
layercard(4.82, "LAYER 2 · 5mC", RGBColor(0x9B,0xE8,0xDC), "The software",
          [[("DNA methylation — reversible marks that ", {"size": 14.5, "color": ICE}),
            ("switch genes up or down", {"size": 14.5, "bold": True, "color": WHITE}),
            (".", {"size": 14.5, "color": ICE})],
           [("", {"size": 6})],
           [("Changes ", {"size": 14.5, "color": ICE}),
            ("how genes are used.", {"size": 14.5, "bold": True, "color": WHITE})]])
layercard(8.74, "LAYER 3 · ANNOTATION", RGBColor(0xF1,0xC4,0x8C), "The interpreter",
          [[("RefSeq function turns coordinates into ", {"size": 14.5, "color": ICE}),
            ("gene names, products, GO", {"size": 14.5, "bold": True, "color": WHITE}),
            (".", {"size": 14.5, "color": ICE})],
           [("", {"size": 6})],
           [("Tells us ", {"size": 14.5, "color": ICE}),
            ("which genes & why.", {"size": 14.5, "bold": True, "color": WHITE})]])

txt(s, 0.9, 6.25, 11.5, 0.6,
    [[("Same PacBio HiFi reads resolve PAV and methylation at once; annotation joins them to biology.",
       {"size": 15.5, "italic": True, "bold": True, "color": ICE})]], align=PP_ALIGN.CENTER)
page_no(s, num())

# ============================================================ 4. STUDY DESIGN
s = slide(WHITE)
kicker(s, 0.9, 0.6, "Study design & data")
txt(s, 0.9, 1.0, 11.5, 0.9, [[("Eight fish, one platform, three readouts",
                               {"font": HEAD_FONT, "size": 32, "bold": True, "color": INK})]])

# stat tiles
def tile(x, big, small, color):
    rect(s, x, 2.2, 2.75, 1.7, LIGHTBG, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.2, 0.12, 1.7, color)
    txt(s, x+0.32, 2.42, 2.3, 0.9, [[(big, {"font": HEAD_FONT, "size": 38, "bold": True, "color": color})]])
    txt(s, x+0.34, 3.32, 2.3, 0.5, [[(small, {"size": 13, "color": MUTE})]])

tile(0.9, "8", "PacBio HiFi individuals", MID)
tile(3.95, "4 + 4", "lean + siscowet", TEAL)
tile(7.0, "1", "reference: SaNama_1.0", SEAFOAM)
tile(10.05, "46,359", "annotated genes", GOLD)

rect(s, 0.9, 4.25, 11.5, 1.75, LIGHTBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 1.25, 4.5, 11, 0.45, [[("Why HiFi matters here",
                               {"font": HEAD_FONT, "size": 18, "bold": True, "color": INK})]])
bullets(s, 1.25, 5.0, 11, 1.0,
        [[("Long, accurate reads call ", {"size": 15}), ("structural variants (PAV)", {"size": 15, "bold": True}),
          (" and ", {"size": 15}), ("5mC methylation", {"size": 15, "bold": True}),
          (" from the very same molecules — no separate assays, no batch confound.", {"size": 15})]],
        gap=6)
txt(s, 0.9, 6.35, 11.5, 0.6,
    [[("Samples — Lean: bc2041/2068/2069/2070   ·   Siscowet: bc2071/2072/2073/2096   ·   Supporting layer: liver RNA-seq (202 DETs, separate parasite study)",
       {"size": 12, "color": MUTE})]])
page_no(s, num())

# ============================================================ 5. PAV — what it is
s = slide(WHITE)
kicker(s, 0.9, 0.6, "Layer 1 · Presence-Absence Variation")
txt(s, 0.9, 1.0, 7.5, 1.6, [[("Whole segments of DNA present in one ecotype, absent in the other",
                              {"font": HEAD_FONT, "size": 29, "bold": True, "color": INK})]])
bullets(s, 0.9, 3.0, 7.2, 3.0,
        [[("Detected two ways from HiFi alignments:", {"size": 17, "bold": True})],
         [("Deletions", {"size": 16, "bold": True, "color": MID}), (" — coverage drops to zero over a region", {"size": 16})],
         [("Insertions", {"size": 16, "bold": True, "color": CORAL}), (" — novel sequence flagged in read CIGAR strings", {"size": 16})],
         [("Classified as lean-specific, siscowet-specific, or shared", {"size": 16})]],
        gap=12)
# side visual: two strands, one with a gap
vx = 8.7
rect(s, vx, 2.4, 3.7, 3.6, LIGHTBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, vx, 2.55, 3.7, 0.4, [[("Lean", {"size": 14, "bold": True, "color": TEAL})]], align=PP_ALIGN.CENTER)
for i,(seg,col) in enumerate([(0,MID),(1,MID),(2,CORAL),(3,MID)]):
    rect(s, vx+0.35+i*0.78, 3.0, 0.68, 0.45, col)
txt(s, vx, 4.05, 3.7, 0.4, [[("Siscowet", {"size": 14, "bold": True, "color": MID})]], align=PP_ALIGN.CENTER)
for i,(seg,col,present) in enumerate([(0,MID,1),(1,RGBColor(0xD7,0xDE,0xE2),0),(2,CORAL,1),(3,MID,1)]):
    rect(s, vx+0.35+i*0.78, 4.45, 0.68, 0.45, col)
txt(s, vx+0.35+0.78, 4.95, 0.9, 0.5, [[("deletion", {"size": 11, "italic": True, "color": MUTE})]], align=PP_ALIGN.CENTER)
txt(s, vx, 5.5, 3.7, 0.4, [[("same locus → different content", {"size": 12, "italic": True, "color": MUTE})]], align=PP_ALIGN.CENTER)
page_no(s, num())

# ============================================================ 6. PAV results (chart)
s = slide(WHITE)
kicker(s, 0.9, 0.6, "Layer 1 · Results")
txt(s, 0.9, 1.0, 11.5, 0.9, [[("Siscowet carries more ecotype-specific structure",
                               {"font": HEAD_FONT, "size": 30, "bold": True, "color": INK})]])

cd = CategoryChartData()
cd.categories = ["Lean-specific", "Siscowet-specific", "Shared"]
cd.add_series("Insertions", (770891, 1086799, 0))
cd.add_series("Deletions", (225337, 245906, 0))
cd.add_series("Shared variants", (0, 0, 878372))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,
                            Inches(0.9), Inches(2.05), Inches(7.4), Inches(4.0), cd)
chart = gframe.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(12)
plot = chart.plots[0]
plot.gap_width = 80
cols = [CORAL, MID, TEAL]
for ser, c in zip(chart.series, cols):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = c
cat = chart.category_axis
cat.tick_labels.font.size = Pt(12)
cat.tick_labels.font.bold = True
val = chart.value_axis
val.tick_labels.font.size = Pt(10)
val.has_major_gridlines = True

# big number callouts on right
def callout(y, num_s, label, color):
    txt(s, 8.7, y, 3.9, 0.7, [[(num_s, {"font": HEAD_FONT, "size": 30, "bold": True, "color": color})]])
    txt(s, 8.72, y+0.58, 3.9, 0.4, [[(label, {"size": 12.5, "color": MUTE})]])
callout(2.1, "996,228", "lean-specific variants", CORAL)
callout(3.25, "1,332,705", "siscowet-specific variants", MID)
callout(4.4, "878,372", "shared between ecotypes", TEAL)
txt(s, 8.7, 5.4, 3.9, 0.9, [[("≈ 34% more ecotype-specific variation in siscowet — but read with reference bias in mind",
                              {"size": 13, "italic": True, "bold": True, "color": INK})]])
txt(s, 0.9, 6.25, 11.5, 0.6,
    [[("Interpretation focuses on ", {"size": 13, "color": MUTE}),
      ("3,465 stringent siscowet-specific deletions", {"size": 13, "bold": True, "color": MID}),
      (" (>100 bp, all-4-vs-none) — the high-confidence set carried into annotation.", {"size": 13, "color": MUTE})]])
page_no(s, num())

# ============================================================ 7. PAV mechanism
s = slide(LIGHTBG)
kicker(s, 0.9, 0.6, "Layer 1 · From variants to phenotype")
txt(s, 0.9, 1.0, 11.5, 0.9, [[("How presence-absence becomes ecotype",
                               {"font": HEAD_FONT, "size": 32, "bold": True, "color": INK})]])
def mech(x, title, body, color, icon):
    rect(s, x, 2.2, 3.65, 3.4, WHITE, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x+0.35, 2.55, 0.7, 0.7, color, shape=MSO_SHAPE.OVAL)
    txt(s, x+0.35, 2.62, 0.7, 0.55, [[(icon, {"size": 24, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER)
    txt(s, x+0.35, 3.45, 3.0, 0.6, [[(title, {"font": HEAD_FONT, "size": 19, "bold": True, "color": INK})]])
    txt(s, x+0.35, 4.15, 3.0, 1.3, [[(body, {"size": 14.5, "color": INK})]])
mech(0.9, "Add / remove genes", "A PAV overlapping a coding region adds or deletes a gene outright — a present-vs-absent function.", MID, "±")
mech(4.85, "Shift gene dosage", "Copy-number change tunes expression up or down without touching the coding sequence.", TEAL, "×")
mech(8.8, "Rewire regulation", "Indels in promoters / enhancers reposition regulatory elements, altering control.", CORAL, "~")
txt(s, 0.9, 5.95, 11.5, 1.0,
    [[("Candidate targets for ecotype divergence: ", {"size": 16, "color": INK}),
      ("lipid metabolism, depth & pressure tolerance, sensory adaptation", {"size": 16, "bold": True, "color": MID}),
      (" — annotation (Layer 3) tells us which genes are actually hit.", {"size": 16, "color": INK})]])
page_no(s, num())

# ============================================================ 8. Methylation intro
s = slide(WHITE)
kicker(s, 0.9, 0.6, "Layer 2 · DNA Methylation", SEAFOAM)
txt(s, 0.9, 1.0, 7.6, 1.6, [[("Same sequence, different settings",
                              {"font": HEAD_FONT, "size": 31, "bold": True, "color": INK})]])
bullets(s, 0.9, 2.9, 7.3, 3.2,
        [[("5-methylcytosine (5mC) at CpG sites", {"size": 16, "bold": True, "color": SEAFOAM}),
          (" — read directly off HiFi data", {"size": 16})],
         [("Methylation typically ", {"size": 16}), ("silences", {"size": 16, "bold": True}),
          (" genes and represses transposable elements", {"size": 16})],
         [("Reversible & potentially heritable", {"size": 16, "bold": True}),
          (" — a route to plastic phenotype", {"size": 16})],
         [("We test each CpG for ecotype differences, then group into regions (DMRs)", {"size": 16})]],
        gap=13, bullet_color=SEAFOAM)
# side: methyl tag concept
vx = 8.8
rect(s, vx, 2.5, 3.6, 3.3, LIGHTBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, vx+0.4, 4.4, 2.8, 0.5, RGBColor(0xC9,0xD2,0xD7))
for i in range(5):
    on = i in (1,3)
    c = SEAFOAM if on else RGBColor(0x9A,0xA8,0xAF)
    cx = vx+0.55+i*0.55
    rect(s, cx, 4.45, 0.16, 0.4, RGBColor(0x6B,0x7A,0x82))  # CpG tick
    if on:
        rect(s, cx-0.12, 3.95, 0.4, 0.4, c, shape=MSO_SHAPE.OVAL)
        txt(s, cx-0.12, 4.0, 0.4, 0.32, [[("m", {"size": 14, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER)
txt(s, vx, 5.15, 3.6, 0.5, [[("methyl marks tune the same DNA", {"size": 12.5, "italic": True, "color": MUTE})]], align=PP_ALIGN.CENTER)
page_no(s, num())

# ============================================================ 9. Methylation results
s = slide(WHITE)
kicker(s, 0.9, 0.6, "Layer 2 · Results", SEAFOAM)
txt(s, 0.9, 1.0, 11.5, 0.9, [[("From half a million sites to 302 regions",
                               {"font": HEAD_FONT, "size": 31, "bold": True, "color": INK})]])
# funnel of three stat tiles
def funnel(x, w, big, label, color):
    rect(s, x, 2.25, w, 1.55, LIGHTBG, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.25, w, 0.14, color)
    txt(s, x, 2.5, w, 0.8, [[(big, {"font": HEAD_FONT, "size": 34, "bold": True, "color": color})]], align=PP_ALIGN.CENTER)
    txt(s, x, 3.32, w, 0.4, [[(label, {"size": 13, "color": MUTE})]], align=PP_ALIGN.CENTER)
funnel(0.9, 3.4, "540,040", "CpG sites tested", TEAL)
txt(s, 4.35, 2.85, 0.55, 0.5, [[("→", {"size": 26, "bold": True, "color": MUTE})]], align=PP_ALIGN.CENTER)
funnel(4.95, 3.4, "4,440", "significant DMCs (p < 0.05)", MID)
txt(s, 8.4, 2.85, 0.55, 0.5, [[("→", {"size": 26, "bold": True, "color": MUTE})]], align=PP_ALIGN.CENTER)
funnel(9.0, 3.4, "302", "differentially methylated regions", SEAFOAM)

# directional split chart
cd = CategoryChartData()
cd.categories = ["Hypermethylated", "Hypomethylated"]
cd.add_series("DMRs in siscowet", (20, 282))
gframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            Inches(0.9), Inches(4.2), Inches(6.3), Inches(2.45), cd)
ch = gframe.chart
ch.has_title = False
ch.has_legend = False
pl = ch.plots[0]
pl.has_data_labels = True
pl.data_labels.font.size = Pt(13)
pl.data_labels.font.bold = True
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
ser = ch.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = SEAFOAM
ch.category_axis.tick_labels.font.size = Pt(12)
ch.category_axis.tick_labels.font.bold = True
ch.value_axis.tick_labels.font.size = Pt(9)

txt(s, 7.5, 4.3, 4.9, 2.4,
    [[("The asymmetry is the story.", {"size": 18, "bold": True, "font": HEAD_FONT, "color": INK})],
     [("", {"size": 5})],
     [("282 of 302 DMRs are ", {"size": 15, "color": INK}),
      ("hypomethylated", {"size": 15, "bold": True, "color": SEAFOAM}),
      (" in siscowet — a genome-wide tilt toward de-repression.", {"size": 15, "color": INK})],
     [("", {"size": 5})],
     [("But 0 single CpGs survive q < 0.1 — we lead with the DMR level, not individual sites.",
       {"size": 13, "italic": True, "color": MUTE})]])
page_no(s, num())

# ============================================================ 10. Methylation mechanism
s = slide(DEEP)
kicker(s, 0.9, 0.6, "Layer 2 · Interpreting the tilt", ICE)
txt(s, 0.9, 1.0, 11.5, 0.9, [[("Why a hypomethylated siscowet genome matters",
                               {"font": HEAD_FONT, "size": 30, "bold": True, "color": WHITE})]])
def dcard(x, title, body):
    rect(s, x, 2.25, 3.65, 3.5, MID, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x+0.35, 2.55, 3.0, 0.9, [[(title, {"font": HEAD_FONT, "size": 19, "bold": True, "color": WHITE})]])
    txt(s, x+0.35, 3.55, 3.0, 2.0, [[(body, {"size": 14.5, "color": ICE})]])
dcard(0.9, "Genes switch on", "Lower methylation near promoters can de-repress genes — potentially the lipid & depth programs that define siscowet.")
dcard(4.85, "Mobile elements wake", "Genome-wide hypomethylation can re-activate transposable elements, a documented engine of rapid divergence.")
dcard(8.8, "A heritable dial", "Methylation state can persist across generations — divergence without waiting for new mutations.")
txt(s, 0.9, 6.1, 11.5, 0.7,
    [[("Two independent signals, same direction: siscowet shows ", {"size": 16, "color": ICE}),
      ("more structural variation and broad hypomethylation.", {"size": 16, "bold": True, "color": WHITE})]],
    align=PP_ALIGN.CENTER)
page_no(s, num())

# ============================================================ 11. Annotation: the join
s = slide(WHITE)
kicker(s, 0.9, 0.6, "Layer 3 · Functional annotation  (new)", GOLD)
txt(s, 0.9, 1.0, 11.5, 0.9, [[("Turning coordinates into candidate genes",
                               {"font": HEAD_FONT, "size": 31, "bold": True, "color": INK})]])

# annotation backbone tiles
def atile(x, big, label, color):
    rect(s, x, 2.15, 3.65, 1.6, LIGHTBG, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.15, 3.65, 0.13, color)
    txt(s, x, 2.4, 3.65, 0.75, [[(big, {"font": HEAD_FONT, "size": 33, "bold": True, "color": color})]], align=PP_ALIGN.CENTER)
    txt(s, x, 3.2, 3.65, 0.45, [[(label, {"size": 12.5, "color": MUTE})]], align=PP_ALIGN.CENTER)
atile(0.9, "46,359", "genes annotated (RefSeq)", GOLD)
atile(4.83, "46,231", "with a product description", TEAL)
atile(8.76, "34,367", "with ≥1 Gene Ontology term", SEAFOAM)

rect(s, 0.9, 4.1, 11.5, 2.05, LIGHTBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 1.25, 4.3, 11, 0.45, [[("The move that makes everything interpretable",
                               {"font": HEAD_FONT, "size": 18, "bold": True, "color": INK})]])
bullets(s, 1.25, 4.85, 10.9, 1.2,
        [[("Strand-aware assignment", {"size": 14.5, "bold": True, "color": GOLD}),
          (" of every DMR / DMC / stringent-PAV to genes (promoter ±2 kb, flank ±5 kb)", {"size": 14.5})],
         [("2,036 candidate genes", {"size": 14.5, "bold": True, "color": GOLD}),
          (" sit within 5 kb of a differential feature — ranked by convergence, placement & expression", {"size": 14.5})],
         [("149/302 DMRs and 1,543/3,465 deletions land near a gene; ", {"size": 14.5}),
          ("54 deletions overlap an exon", {"size": 14.5, "bold": True}),
          (" (candidate copy/LOF)", {"size": 14.5})]],
        gap=8, bullet_color=GOLD)
page_no(s, num())

# ============================================================ 12. Convergent + top candidates
s = slide(WHITE)
kicker(s, 0.9, 0.55, "Layer 3 · Where the layers converge", GOLD)
txt(s, 0.9, 0.95, 11.5, 0.8, [[("Four genes carry both a DMR and a siscowet deletion",
                                {"font": HEAD_FONT, "size": 28, "bold": True, "color": INK})]])

# candidate figure on the right
import os
fig = "/Users/sr320/GitHub/project-lake-trout/figures/18-top-candidates.png"
if os.path.exists(fig):
    s.shapes.add_picture(fig, Inches(7.15), Inches(1.95), width=Inches(5.4))
    txt(s, 7.15, 6.55, 5.4, 0.4,
        [[("Top protein-coding candidates by integrated rank score.",
           {"size": 11, "italic": True, "color": MUTE})]], align=PP_ALIGN.CENTER)

# convergent table (left)
def crow(y, gene, note, color):
    rect(s, 0.9, y, 6.0, 0.62, LIGHTBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.9, y, 0.1, 0.62, color)
    txt(s, 1.15, y+0.07, 3.0, 0.5, [[(gene, {"size": 14, "bold": True, "color": INK})]])
    txt(s, 3.95, y+0.1, 2.85, 0.45, [[(note, {"size": 11.5, "color": MUTE})]])
txt(s, 0.9, 1.9, 6.0, 0.35, [[("4 CONVERGENT LOCI", {"size": 12, "bold": True, "color": GOLD, "spacing": 1.5})]])
crow(2.3, "znf883-like", "exon DMR + exonic deletion · TOP", GOLD)
crow(3.0, "XlCGF57.1-like", "intron hypo-DMR + nearby deletion", TEAL)
crow(3.7, "septin-9-like", "intron hyper-DMR + nearby deletion", TEAL)
crow(4.4, "LOC120039781", "uncharacterized · intron hypo-DMR", MUTE)

txt(s, 0.9, 5.2, 6.0, 0.35, [[("EXONIC DELETIONS IN LIPID GENES", {"size": 12, "bold": True, "color": CORAL, "spacing": 1.2})]])
bullets(s, 0.9, 5.6, 6.0, 1.3,
        [[("angptl5", {"size": 13.5, "bold": True, "color": INK}), ("  ·  ", {"size": 13.5, "color": MUTE}),
          ("mogat2", {"size": 13.5, "bold": True, "color": INK}), ("  ·  ", {"size": 13.5, "color": MUTE}),
          ("epoxide hydrolase 1", {"size": 13.5, "bold": True, "color": INK})],
         [("Lipid handling — consistent with the defining high-lipid siscowet phenotype",
           {"size": 12.5, "italic": True, "color": MUTE})]],
        gap=5, bullet_color=CORAL)
page_no(s, num())

# ============================================================ 13. GO enrichment + phenotype axes
s = slide(LIGHTBG)
kicker(s, 0.9, 0.6, "Layer 3 · Function & phenotype")
txt(s, 0.9, 1.0, 11.5, 0.9, [[("What the candidate set points toward",
                               {"font": HEAD_FONT, "size": 31, "bold": True, "color": INK})]])

# GO enrichment mini-table (left)
rect(s, 0.9, 2.1, 5.55, 4.05, WHITE, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 1.2, 2.3, 5.0, 0.4, [[("GO enrichment — deletion set", {"font": HEAD_FONT, "size": 16, "bold": True, "color": INK})]])
def gorow(y, term, fold, fdr, strong=False):
    tc = INK if strong else MUTE
    txt(s, 1.2, y, 3.3, 0.4, [[(term, {"size": 12.5, "bold": strong, "color": tc})]])
    txt(s, 4.55, y, 0.8, 0.4, [[(fold, {"size": 12.5, "color": tc})]], align=PP_ALIGN.RIGHT)
    txt(s, 5.35, y, 0.95, 0.4, [[(fdr, {"size": 12.5, "color": tc})]], align=PP_ALIGN.RIGHT)
txt(s, 1.2, 2.78, 3.3, 0.35, [[("term", {"size": 10.5, "bold": True, "color": MUTE, "spacing": 1})]])
txt(s, 4.55, 2.78, 0.8, 0.35, [[("FOLD", {"size": 10.5, "bold": True, "color": MUTE})]], align=PP_ALIGN.RIGHT)
txt(s, 5.35, 2.78, 0.95, 0.35, [[("FDR", {"size": 10.5, "bold": True, "color": MUTE})]], align=PP_ALIGN.RIGHT)
gorow(3.2, "Ca²⁺ transmembrane transport", "3.7", "5.6e-4", True)
gorow(3.68, "Neuron projection development", "2.4", "2.6e-3", True)
gorow(4.16, "Calcium channel complex", "4.6", "3.0e-3")
gorow(4.64, "Calcium ion transport", "3.0", "3.0e-3", True)
gorow(5.12, "Lipid / phospholipid binding", "1.5", "ns")
txt(s, 1.2, 5.65, 5.0, 0.45, [[("Calcium transport = most defensible; carries a gene-length caveat. DMR set is a tandem-cluster artifact.",
                                {"size": 10.5, "italic": True, "color": MUTE})]])

# phenotype axes (right)
def axis(x, y, w, icon, head, body, color):
    rect(s, x, y, w, 0.92, WHITE, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, 0.1, 0.92, color)
    txt(s, x+0.28, y+0.13, 0.6, 0.6, [[(icon, {"size": 20})]])
    txt(s, x+0.92, y+0.1, w-1.0, 0.4, [[(head, {"font": HEAD_FONT, "size": 14, "bold": True, "color": INK})]])
    txt(s, x+0.92, y+0.47, w-1.0, 0.4, [[(body, {"size": 11, "color": MUTE})]])
txt(s, 6.75, 2.1, 5.6, 0.35, [[("HYPOTHESIZED PHENOTYPE AXES", {"size": 12, "bold": True, "color": GOLD, "spacing": 1.2})]])
axis(6.75, 2.5, 5.65, "🫧", "Lipid & energy storage", "angptl5, mogat2, ephx1 — high-lipid siscowet / buoyancy", CORAL)
axis(6.75, 3.5, 5.65, "🐠", "Body shape, growth & muscle", "rbm24b + growth GO — elongate-lean vs. robust-siscowet", TEAL)
axis(6.75, 4.5, 5.65, "⚡", "Calcium / sensory-neural", "Ca²⁺ & neuron-projection GO — deep, dark, high-pressure", MID)
axis(6.75, 5.5, 5.65, "🛡️", "Immune & adhesion", "a2m, CEACAM, DMBT1 — some real, some reference-divergent", SEAFOAM)
page_no(s, num())

# ============================================================ 14. Guardrails
s = slide(DEEP)
kicker(s, 0.9, 0.6, "Read this carefully", ICE)
txt(s, 0.9, 1.0, 11.5, 0.9, [[("Hypothesis-generating, not causal",
                               {"font": HEAD_FONT, "size": 31, "bold": True, "color": WHITE})]])
def guard(x, y, w, head, body):
    rect(s, x, y, w, 1.55, MID, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, 0.1, 1.55, CORAL)
    txt(s, x+0.32, y+0.18, w-0.6, 0.45, [[(head, {"font": HEAD_FONT, "size": 15.5, "bold": True, "color": WHITE})]])
    txt(s, x+0.32, y+0.66, w-0.6, 0.8, [[(body, {"size": 12.5, "color": ICE})]])
guard(0.9, 2.15, 5.65, "Lean-background reference",
      "SaNama_1.0 is a doubled-haploid lean fish. Siscowet diverges more, inflating its apparent deletions — counts are not magnitude-comparable.")
guard(6.78, 2.15, 5.65, "No single CpG survives correction",
      "0 DMCs at q < 0.1. Interpretation leads with DMR-level and stringent-PAV sets; single-CpG hits are suggestive only.")
guard(0.9, 3.95, 5.65, "Expression support is weak by design",
      "Liver RNA-seq comes from a separate parasite study with different fish — orthogonal support, never confirmation.")
guard(6.78, 3.95, 5.65, "Enrichment confounders",
      "PAV GO carries a gene-length bias (long Ca²⁺ channels); the DMR GO signal is a single histone/znf883 tandem cluster.")
txt(s, 0.9, 5.95, 11.5, 0.8,
    [[("The deliverable is a ranked, annotated shortlist of candidates — ", {"size": 16, "color": ICE}),
      ("associations on one reference, awaiting functional validation.", {"size": 16, "bold": True, "color": WHITE})]],
    align=PP_ALIGN.CENTER)
page_no(s, num())

# ============================================================ 15. Takeaways
s = slide(LIGHTBG)
kicker(s, 0.9, 0.6, "Take-home")
txt(s, 0.9, 1.0, 11.5, 0.9, [[("Three things to remember",
                               {"font": HEAD_FONT, "size": 33, "bold": True, "color": INK})]])
def take(y, n, head, body, color):
    rect(s, 0.9, y, 11.5, 1.45, WHITE, shadow=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.9, y, 0.14, 1.45, color)
    rect(s, 1.25, y+0.32, 0.8, 0.8, color, shape=MSO_SHAPE.OVAL)
    txt(s, 1.25, y+0.4, 0.8, 0.6, [[(n, {"font": HEAD_FONT, "size": 26, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER)
    txt(s, 2.4, y+0.22, 9.6, 0.5, [[(head, {"font": HEAD_FONT, "size": 19, "bold": True, "color": INK})]])
    txt(s, 2.4, y+0.78, 9.6, 0.55, [[(body, {"size": 14.5, "color": MUTE})]])
take(2.15, "1", "Divergence is written at two levels, read at a third",
     "Structure (what genes exist) + methylation (how they're used), made interpretable by annotation.", MID)
take(3.78, "2", "Annotation produced a ranked shortlist",
     "2,036 candidates, 4 convergent loci led by znf883-like, lipid-gene deletions, and calcium-transport GO.", GOLD)
take(5.41, "3", "Read it as hypotheses, not mechanisms",
     "Lean-reference bias, no FDR-significant CpG, and weak expression support keep every link associative.", CORAL)
page_no(s, num())

# ============================================================ 16. Closing
s = slide(DEEP)
rect(s, 0, 5.4, 13.333, 2.1, MID)
rect(s, 0, 6.4, 13.333, 1.1, TEAL)
rect(s, 0, 7.05, 13.333, 0.45, SEAFOAM)
txt(s, 0.9, 1.6, 11.5, 1.6, [
    [("From variation to candidate genes", {"font": HEAD_FONT, "size": 38, "bold": True, "color": WHITE})],
])
txt(s, 0.92, 2.95, 11, 1.5,
    [[("Methylation and structural variation each track the lean–siscowet split; ",
       {"size": 19, "color": ICE})],
     [("functional annotation turns where they converge into the genes that may make an ecotype.",
       {"size": 19, "bold": True, "color": WHITE})]], space_after=4)
txt(s, 0.92, 5.65, 11.5, 0.5,
    [[("Explore the data — interactive genome browsers (IGV.js + JBrowse 2):", {"size": 14, "color": ICE})]])
txt(s, 0.92, 6.55, 11.5, 0.5,
    [[("sr320.github.io/project-lake-trout", {"size": 16, "bold": True, "color": WHITE})]])
txt(s, 0.92, 7.12, 9, 0.35, [[("Goetz · White · Gallardo-Escarate · Roberts  ·  Roberts Lab, UW SAFS",
                               {"size": 11.5, "color": ICE})]])

out = "/Users/sr320/GitHub/project-lake-trout/output/deck/lake-trout-ecotypes-PAV-methylation.pptx"
prs.save(out)
print("Saved", out, "slides:", len(prs.slides._sldIdLst))
